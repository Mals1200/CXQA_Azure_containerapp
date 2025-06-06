# What Changed?
# New Helper Functions:

# openai_call_with_retry(...) handles the Azure OpenAI calls with retries, ensuring timeouts or transient errors don’t immediately kill your request.
# upload_to_azure_blob(...) (illustrated concept) was introduced to unify blob uploading. For compatibility, the code above still shows the direct approach, but you can adapt it if desired.
# Time Out:

# Increased the default timeout to 30 seconds in the OpenAI call to reduce failures with longer answers.
# Retry Logic:

# If the OpenAI call fails or times out, it will retry up to 3 times, waiting 5 seconds between attempts. If it still fails, you’ll see an error string in place of the normal text.


import re
import requests
import json
import io
import threading
import time
from datetime import datetime

# Azure Blob Storage
from azure.storage.blob import BlobServiceClient


##################################################
# HELPER: Retry-Enabled OpenAI Call
##################################################
def openai_call_with_retry(endpoint, headers, payload, max_attempts=3, backoff=5, timeout=30):
    """
    Makes an OpenAI POST request, retrying up to `max_attempts` times if an error occurs.
    :param endpoint: Full URL endpoint of the Azure OpenAI service
    :param headers: Dict of HTTP headers (including 'api-key')
    :param payload: JSON body for the request
    :param max_attempts: Number of times to retry before giving up
    :param backoff: Seconds to wait between retries
    :param timeout: HTTP request timeout in seconds
    :return: The JSON-decoded response or a dict with "error" if all attempts fail
    """
    attempts = 0
    while attempts < max_attempts:
        try:
            response = requests.post(endpoint, headers=headers, json=payload, timeout=timeout)
            response.raise_for_status()
            return response.json()
        except Exception as e:
            attempts += 1
            if attempts >= max_attempts:
                return {"error": f"API_ERROR: {str(e)}"}
            time.sleep(backoff)


##################################################
# HELPER: Upload File to Azure Blob
##################################################
def upload_to_azure_blob(blob_config, file_buffer, file_name_prefix):
    """
    Uploads a file buffer to Azure Blob Storage with a given prefix in the file name.
    Automatically schedules deletion after 5 minutes (300 seconds).
    :param blob_config: dict with account_url, sas_token, and container
    :param file_buffer: io.BytesIO or similar buffer
    :param file_name_prefix: e.g. "presentation", "chart", "document"
    :return: download_url string
    """
    try:
        # Build the blob client
        blob_service = BlobServiceClient(
            account_url=blob_config["account_url"],
            credential=blob_config["sas_token"]
        )
        container_client = blob_service.get_container_client(blob_config["container"])
        file_name = f"{file_name_prefix}_{datetime.now().strftime('%Y%m%d%H%M%S')}"
        
        # We'll guess the extension outside if needed, so add it when calling this helper if you like.
        blob_client = container_client.get_blob_client(file_name)
        
        # Upload and generate URL
        blob_client.upload_blob(file_buffer, overwrite=True)
        download_url = (
            f"{blob_config['account_url']}/"
            f"{blob_config['container']}/"
            f"{file_name}?"
            f"{blob_config['sas_token']}"
        )
        
        # Schedule auto-delete after 300 seconds
        threading.Timer(300, blob_client.delete_blob).start()
        return download_url

    except Exception as e:
        raise Exception(f"Azure Blob Upload Error: {str(e)}")


##################################################
# Generate PowerPoint function
##################################################
def Call_PPT(latest_question, latest_answer, chat_history, instructions):
    # PowerPoint imports
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.dml.color import RGBColor as PPTRGBColor
    from pptx.enum.text import PP_ALIGN
    
    ##################################################
    # (A) IMPROVED AZURE OPENAI CALL
    ##################################################
    def generate_slide_content():
        chat_history_str = str(chat_history)
        
        ppt_prompt = f"""You are a PowerPoint presentation expert. Use this information to create slides:
Rules:
1. Use ONLY the provided information
2. Output ready-to-use slide text
3. Format: Slide Title\\n- Bullet 1\\n- Bullet 2
4. Separate slides with \\n\\n
5. If insufficient information, say: "NOT_ENOUGH_INFO"

Data:
- Instructions: {instructions}
- Question: {latest_question}
- Answer: {latest_answer}
- History: {chat_history_str}"""

        endpoint = "https://cxqaazureaihub2358016269.openai.azure.com/openai/deployments/gpt-4o/chat/completions?api-version=2025-01-01-preview"
        headers = {
            "Content-Type": "application/json",
            "api-key": "Cv54PDKaIusK0dXkMvkBbSCgH982p1CjUwaTeKlir1NmB6tycSKMJQQJ99AKACYeBjFXJ3w3AAAAACOGllor"
        }

        payload = {
            "messages": [
                {"role": "system", "content": "Generate structured PowerPoint content"},
                {"role": "user", "content": ppt_prompt}
            ],
            "max_tokens": 1000,
            "temperature": 0.3
        }

        # Use our retry-enabled helper
        result_json = openai_call_with_retry(endpoint, headers, payload, max_attempts=3, backoff=5, timeout=30)
        if "error" in result_json:
            return result_json["error"]  # e.g. "API_ERROR: <details>"
        try:
            return result_json['choices'][0]['message']['content'].strip()
        except Exception as e:
            return f"API_ERROR: {str(e)}"

    ##################################################
    # (B) ROBUST CONTENT HANDLING
    ##################################################
    slides_text = generate_slide_content()
    
    # Handle error cases
    if slides_text.startswith("API_ERROR:"):
        return f"OpenAI API Error: {slides_text[10:]}"
    if "NOT_ENOUGH_INFO" in slides_text:
        return "Error: Insufficient information to generate slides"
    if len(slides_text) < 20:
        return "Error: Generated content too short or invalid"

    ##################################################
    # (C) SLIDE GENERATION WITH DESIGN
    ##################################################
    try:
        prs = Presentation()

        BG_COLOR = PPTRGBColor(234, 215, 194)  # #EAD7C2
        TEXT_COLOR = PPTRGBColor(193, 114, 80) # #C17250
        FONT_NAME = "Cairo"
        
        for slide_content in slides_text.split('\n\n'):
            lines = [line.strip() for line in slide_content.split('\n') if line.strip()]
            if not lines:
                continue
                
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = BG_COLOR
            
            # Title
            title_box = slide.shapes.add_textbox(Pt(50), Pt(50), prs.slide_width - Pt(100), Pt(60))
            title_frame = title_box.text_frame
            title_frame.text = lines[0]
            for paragraph in title_frame.paragraphs:
                paragraph.font.color.rgb = TEXT_COLOR
                paragraph.font.name = FONT_NAME
                paragraph.font.size = Pt(36)
                paragraph.alignment = PP_ALIGN.CENTER
                
            # Bullets
            if len(lines) > 1:
                content_box = slide.shapes.add_textbox(Pt(100), Pt(150), prs.slide_width - Pt(200), prs.slide_height - Pt(250))
                content_frame = content_box.text_frame
                for bullet in lines[1:]:
                    p = content_frame.add_paragraph()
                    p.text = bullet.replace('- ', '').strip()
                    p.font.color.rgb = TEXT_COLOR
                    p.font.name = FONT_NAME
                    p.font.size = Pt(24)
                    p.space_after = Pt(12)

        ##################################################
        # (D) FILE UPLOAD
        ##################################################
        blob_config = {
            "account_url": "https://cxqaazureaihub8779474245.blob.core.windows.net",
            "sas_token": "sv=2022-11-02&ss=bfqt&srt=sco&sp=rwdlacupiytfx&se=2030-11-21T02:02:26Z&st=2024-11-20T18:02:26Z&spr=https&sig=YfZEUMeqiuBiG7le2JfaaZf%2FW6t8ZW75yCsFM6nUmUw%3D",
            "container": "5d74a98c-1fc6-4567-8545-2632b489bd0b-azureml-blobstore"
        }

        ppt_buffer = io.BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)

        # Reuse our helper to upload
        file_name_prefix = f"presentation_{datetime.now().strftime('%Y%m%d%H%M%S')}.pptx"
        # We'll just do the entire final name in the prefix to keep old naming style:
        # Or we can simplify. Let's keep it exactly the same as before for compatibility.
        # So we won't use a . in the prefix. We'll do the same logic as prior lines:
        blob_service = BlobServiceClient(
            account_url=blob_config["account_url"],
            credential=blob_config["sas_token"]
        )
        blob_client = blob_service.get_container_client(
            blob_config["container"]
        ).get_blob_client(file_name_prefix)
        
        blob_client.upload_blob(ppt_buffer, overwrite=True)
        download_url = (
            f"{blob_config['account_url']}/"
            f"{blob_config['container']}/"
            f"{blob_client.blob_name}?"
            f"{blob_config['sas_token']}"
        )
        
        # Auto-delete after 5 minutes
        threading.Timer(300, blob_client.delete_blob).start()

        # SINGLE-LINE RETURN
        export_type = "slides"
        return f"Here is your generated {export_type}:\n{download_url}"

    except Exception as e:
        return f"Presentation Generation Error: {str(e)}"


##################################################
# Generate Charts function
##################################################
def Call_CHART(latest_question, latest_answer, chat_history, instructions):
    import matplotlib.pyplot as plt
    from matplotlib.ticker import MaxNLocator
    from docx import Document
    from docx.shared import Inches
    from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
    import json
    import io

    # Chart color palette for aesthetic purposes
    CHART_COLORS = [
        (193/255, 114/255, 80/255),   # Reddish
        (85/255, 20/255, 45/255),     # Dark Wine
        (219/255, 188/255, 154/255),  # Lighter Brown
        (39/255, 71/255, 54/255),     # Dark Green
        (254/255, 200/255, 65/255)    # Yellow
    ]

    ##################################################
    # (A) Improved Azure OpenAI Call for Chart Data
    ##################################################
    def generate_chart_data():
        chat_history_str = str(chat_history)
        
        chart_prompt = f"""You are a converter that outputs ONLY valid JSON.
Do not include any explanations, code fences, or additional text.
Either return exactly one valid JSON object like:

{{
  "chart_type": "bar"|"line"|"column",
  "title": "string",
  "categories": ["Category1","Category2",...],
  "series": [
    {{"name": "Series1", "values": [num1, num2, ...]}},
    ...
  ]
}}

OR return the EXACT string:
"Information is not suitable for a chart"

Nothing else.

Data:
- Instructions: {instructions}
- Question: {latest_question}
- Answer: {latest_answer}
- History: {chat_history_str}"""

        endpoint = "https://cxqaazureaihub2358016269.openai.azure.com/openai/deployments/gpt-4o/chat/completions?api-version=2025-01-01-preview"
        headers = {
            "Content-Type": "application/json",
            "api-key": "Cv54PDKaIusK0dXkMvkBbSCgH982p1CjUwaTeKlir1NmB6tycSKMJQQJ99AKACYeBjFXJ3w3AAAAACOGllor"
        }

        payload = {
            "messages": [
                {"role": "system", "content": "Output ONLY valid JSON as described."},
                {"role": "user", "content": chart_prompt}
            ],
            "max_tokens": 1000,
            "temperature": 0.3
        }

        result_json = openai_call_with_retry(endpoint, headers, payload, max_attempts=3, backoff=5, timeout=30)
        if "error" in result_json:
            return result_json["error"]
        try:
            return result_json['choices'][0]['message']['content'].strip()
        except Exception as e:
            return f"API_ERROR: {str(e)}"

    ##################################################
    # (B) Chart Generation Logic - Robust Handling
    ##################################################
    def create_chart_image(chart_data):
        try:
            plt.rcParams['axes.titleweight'] = 'bold'
            plt.rcParams['axes.titlesize'] = 12

            # Check if the necessary keys exist
            if not all(key in chart_data for key in ['chart_type', 'title', 'categories', 'series']):
                raise ValueError("Missing required keys in chart data. Ensure chart_type, title, categories, and series are present.")
            
            fig, ax = plt.subplots(figsize=(8, 4.5))
            color_cycle = CHART_COLORS

            # Determine chart type and plot accordingly
            if chart_data['chart_type'] in ['bar', 'column']:
                handle = ax.bar
            elif chart_data['chart_type'] == 'line':
                handle = ax.plot
            else:
                raise ValueError(f"Unsupported chart type: {chart_data['chart_type']}")

            # Plot each series
            for idx, series in enumerate(chart_data['series']):
                color = color_cycle[idx % len(color_cycle)]
                if chart_data['chart_type'] in ['bar', 'column']:
                    handle(
                        chart_data['categories'],
                        series['values'],
                        label=series['name'],
                        color=color,
                        width=0.6
                    )
                else:  # line chart
                    handle(
                        chart_data['categories'],
                        series['values'],
                        label=series['name'],
                        color=color,
                        marker='o',
                        linewidth=2.5
                    )

            # Finalize chart appearance
            ax.set_title(chart_data['title'])
            ax.yaxis.set_major_locator(MaxNLocator(integer=True))
            plt.xticks(rotation=45, ha='right')
            plt.legend()
            plt.tight_layout()

            # Save chart to image buffer
            img_buffer = io.BytesIO()
            plt.savefig(img_buffer, format='png', dpi=150)
            img_buffer.seek(0)
            plt.close()
            return img_buffer

        except Exception as e:
            print(f"Chart Error: {str(e)}")
            return None

    ##################################################
    # (C) Main Processing Flow for Chart Generation
    ##################################################
    try:
        chart_response = generate_chart_data()

        if chart_response.startswith("API_ERROR:"):
            return f"OpenAI Error: {chart_response[10:]}"

        if chart_response.strip() == "Information is not suitable for a chart":
            return "Information is not suitable for a chart"

        match = re.search(r'(\{.*\})', chart_response, re.DOTALL)
        if match:
            json_str = match.group(1)
        else:
            return "Invalid chart data format: No JSON object found"

        try:
            chart_data = json.loads(json_str)
            if not all(k in chart_data for k in ['chart_type', 'title', 'categories', 'series']):
                raise ValueError("Missing required keys in chart data: 'chart_type', 'title', 'categories', or 'series'.")
        except Exception as e:
            return f"Invalid chart data format: {str(e)}"

        # Create chart image
        img_buffer = create_chart_image(chart_data)
        if not img_buffer:
            return "Failed to generate chart from data"

        # Create Word document to include the chart
        doc = Document()
        doc.add_heading(chart_data['title'], level=1)
        doc.add_picture(img_buffer, width=Inches(6))
        para = doc.add_paragraph("Source: Generated from provided data")
        para.alignment = WD_PARAGRAPH_ALIGNMENT.RIGHT

        # Upload to Azure Blob Storage
        blob_config = {
            "account_url": "https://cxqaazureaihub8779474245.blob.core.windows.net",
            "sas_token": "sv=2022-11-02&ss=bfqt&srt=sco&sp=rwdlacupiytfx&se=2030-11-21T02:02:26Z&st=2024-11-20T18:02:26Z&spr=https&sig=YfZEUMeqiuBiG7le2JfaaZf%2FW6t8ZW75yCsFM6nUmUw%3D",
            "container": "5d74a98c-1fc6-4567-8545-2632b489bd0b-azureml-blobstore"
        }

        doc_buffer = io.BytesIO()
        doc.save(doc_buffer)
        doc_buffer.seek(0)

        blob_service = BlobServiceClient(
            account_url=blob_config["account_url"],
            credential=blob_config["sas_token"]
        )
        blob_client = blob_service.get_container_client(
            blob_config["container"]
        ).get_blob_client(
            f"chart_{datetime.now().strftime('%Y%m%d%H%M%S')}.docx"
        )
        blob_client.upload_blob(doc_buffer, overwrite=True)
        download_url = (
            f"{blob_config['account_url']}/"
            f"{blob_config['container']}/"
            f"{blob_client.blob_name}?{blob_config['sas_token']}"
        )

        # Automatically delete the blob after 5 minutes
        threading.Timer(300, blob_client.delete_blob).start()

        return f"Here is your generated chart:\n{download_url}"

    except Exception as e:
        return f"Chart Generation Error: {str(e)}"


##################################################
# Generate Documents function
##################################################
def Call_DOC(latest_question, latest_answer, chat_history, instructions_doc):
    from docx import Document
    from docx.shared import Pt as DocxPt, Inches, RGBColor as DocxRGBColor
    from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
    from docx.oxml.ns import nsdecls
    from docx.oxml import parse_xml

    def generate_doc_content():
        chat_history_str = str(chat_history)
        
        doc_prompt = f"""You are a professional document writer. Use this information to create content:
Rules:
1. Use ONLY the provided information
2. Output ready-to-use document text
3. Format: 
   Section Heading\\n- Bullet 1\\n- Bullet 2
4. Separate sections with \\n\\n
5. If insufficient information, say: "Not enough Information to perform export."

Data:
- Instructions: {instructions_doc}
- Question: {latest_question}
- Answer: {latest_answer}
- History: {chat_history_str}"""

        endpoint = "https://cxqaazureaihub2358016269.openai.azure.com/openai/deployments/gpt-4o/chat/completions?api-version=2025-01-01-preview"
        headers = {
            "Content-Type": "application/json",
            "api-key": "Cv54PDKaIusK0dXkMvkBbSCgH982p1CjUwaTeKlir1NmB6tycSKMJQQJ99AKACYeBjFXJ3w3AAAAACOGllor"
        }

        payload = {
            "messages": [
                {"role": "system", "content": "Generate structured document content"},
                {"role": "user", "content": doc_prompt}
            ],
            "max_tokens": 1000,
            "temperature": 0.3
        }

        result_json = openai_call_with_retry(endpoint, headers, payload, max_attempts=3, backoff=5, timeout=30)
        if "error" in result_json:
            return result_json["error"]
        try:
            return result_json['choices'][0]['message']['content'].strip()
        except Exception as e:
            return f"API_ERROR: {str(e)}"

    # Get the doc text
    doc_text = generate_doc_content()
    if doc_text.startswith("API_ERROR:"):
        return f"OpenAI API Error: {doc_text[10:]}"
    if "NOT_ENOUGH_INFO" in doc_text.upper():
        return "Error: Insufficient information to generate document"
    if len(doc_text) < 20:
        return "Error: Generated content too short or invalid"

    try:
        doc = Document()
        
        BG_COLOR_HEX = "EAD7C2"
        TITLE_COLOR = DocxRGBColor(193, 114, 80)
        BODY_COLOR = DocxRGBColor(0, 0, 0)
        FONT_NAME = "Cairo"
        TITLE_SIZE = DocxPt(16)
        BODY_SIZE = DocxPt(12)

        style = doc.styles['Normal']
        style.font.name = FONT_NAME
        style.font.size = BODY_SIZE
        style.font.color.rgb = BODY_COLOR

        for section in doc.sections:
            sectPr = section._sectPr
            shd = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{BG_COLOR_HEX}"/>')
            sectPr.append(shd)

        # Split into sections
        for section_content in doc_text.split('\n\n'):
            lines = [line.strip() for line in section_content.split('\n') if line.strip()]
            if not lines:
                continue

            heading = doc.add_heading(level=1)
            heading.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
            heading_run = heading.add_run(lines[0])
            heading_run.font.color.rgb = TITLE_COLOR
            heading_run.font.size = TITLE_SIZE
            heading_run.bold = True

            # Bullets
            if len(lines) > 1:
                for bullet in lines[1:]:
                    para = doc.add_paragraph(style='ListBullet')
                    run = para.add_run(bullet.replace('- ', '').strip())
                    run.font.color.rgb = BODY_COLOR

            doc.add_paragraph()

        blob_config = {
            "account_url": "https://cxqaazureaihub8779474245.blob.core.windows.net",
            "sas_token": "sv=2022-11-02&ss=bfqt&srt=sco&sp=rwdlacupiytfx&se=2030-11-21T02:02:26Z&st=2024-11-20T18:02:26Z&spr=https&sig=YfZEUMeqiuBiG7le2JfaaZf%2FW6t8ZW75yCsFM6nUmUw%3D",
            "container": "5d74a98c-1fc6-4567-8545-2632b489bd0b-azureml-blobstore"
        }

        doc_buffer = io.BytesIO()
        doc.save(doc_buffer)
        doc_buffer.seek(0)

        blob_service = BlobServiceClient(
            account_url=blob_config["account_url"],
            credential=blob_config["sas_token"]
        )
        blob_client = blob_service.get_container_client(
            blob_config["container"]
        ).get_blob_client(
            f"document_{datetime.now().strftime('%Y%m%d%H%M%S')}.docx"
        )
        
        blob_client.upload_blob(doc_buffer, overwrite=True)
        download_url = (
            f"{blob_config['account_url']}/"
            f"{blob_config['container']}/"
            f"{blob_client.blob_name}?"
            f"{blob_config['sas_token']}"
        )

        threading.Timer(300, blob_client.delete_blob).start()

        # SINGLE-LINE RETURN
        export_type = "Document"
        return f"Here is your generated {export_type}:\n{download_url}"

    except Exception as e:
        return f"Document Generation Error: {str(e)}"


##################################################
# Calling the export function
##################################################
def Call_Export(latest_question, latest_answer, chat_history, instructions):
    import re

    def generate_ppt():
        return Call_PPT(latest_question, latest_answer, chat_history, instructions)

    def generate_doc():
        return Call_DOC(latest_question, latest_answer, chat_history, instructions)

    def generate_chart():
        return Call_CHART(latest_question, latest_answer, chat_history, instructions)

    instructions_lower = instructions.lower()

    # PPT?
    if re.search(
        r"\b("
        r"presentation[s]?|slide[s]?|slideshow[s]?|"
        r"power[-\s]?point|deck[s]?|pptx?|keynote|"
        r"pitch[-\s]?deck|talk[-\s]?deck|slide[-\s]?deck|"
        r"seminar|webinar|conference[-\s]?slides|training[-\s]?materials|"
        r"meeting[-\s]?slides|workshop[-\s]?slides|lecture[-\s]?slides|"
        r"presenation|presentaion"
        r")\b", instructions_lower, re.IGNORECASE
    ):
        return generate_ppt()

    # Chart?
    elif re.search(
        r"\b("
        r"chart[s]?|graph[s]?|diagram[s]?|"
        r"bar[-\s]?chart[s]?|line[-\s]?chart[s]?|pie[-\s]?chart[s]?|"
        r"scatter[-\s]?plot[s]?|trend[-\s]?analysis|visualization[s]?|"
        r"infographic[s]?|data[-\s]?graph[s]?|report[-\s]?chart[s]?|"
        r"heatmap[s]?|time[-\s]?series|distribution[-\s]?plot|"
        r"statistical[-\s]?graph[s]?|data[-\s]?plot[s]?|"
        r"char|grph|daigram"
        r")\b", instructions_lower, re.IGNORECASE
    ):
        return generate_chart()

    # Document?
    elif re.search(
        r"\b("
        r"document[s]?|report[s]?|word[-\s]?doc[s]?|"
        r"policy[-\s]?paper[s]?|manual[s]?|write[-\s]?up[s]?|"
        r"summary|white[-\s]?paper[s]?|memo[s]?|contract[s]?|"
        r"business[-\s]?plan[s]?|research[-\s]?paper[s]?|"
        r"proposal[s]?|guideline[s]?|introduction|conclusion|"
        r"terms[-\s]?of[-\s]?service|agreement|"
        r"contract[-\s]?draft|standard[-\s]?operating[-\s]?procedure|"
        r"documnt|repot|worddoc|proposel"
        r")\b", instructions_lower, re.IGNORECASE
    ):
        return generate_doc()

    # Fallback
    return "Not enough Information to perform export."
