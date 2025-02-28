# ppt_export_agent.py

import os
import io
import json
import re
import requests
import contextlib
from io import BytesIO
from datetime import datetime
import pandas as pd
from azure.storage.blob import BlobServiceClient
from azure.core.credentials import AzureKeyCredential
from pptx import Presentation
from pptx.util import Inches, Pt

######################################
# FAKE KEYS (same as in ask_func.py) #
######################################
FAKE_LLM_ENDPOINT = (
    "https://cxqaazureaihub2358016269.openai.azure.com/"
    "openai/deployments/gpt-4o-3/chat/completions?api-version=2024-08-01-preview"
)
FAKE_LLM_API_KEY = "Cv54PDKaIusK0dXkMvkBbSCgH982p1CjUwaTeKlir1NmB6tycSKMJQQJ99AKACYeBjFXJ3w3AAAAACOGllor"

BLOB_ACCOUNT_URL = "https://cxqaazureaihub8779474245.blob.core.windows.net"
BLOB_SAS_TOKEN = (
    "sv=2022-11-02&ss=bfqt&srt=sco&sp=rwdlacupiytfx&"
    "se=2030-11-21T02:02:26Z&st=2024-11-20T18:02:26Z&"
    "spr=https&sig=YfZEUMeqiuBiG7le2JfaaZf%2FW6t8ZW75yCsFM6nUmUw%3D"
)
BLOB_CONTAINER_NAME = "5d74a98c-1fc6-4567-8545-2632b489bd0b-azureml-blobstore"
BLOB_FOLDER_PATH = "UI/2024-11-20_142337_UTC/cxqa_data/ppt_exports/"  # choose any subfolder name here if you wish

#######################################
# 1) Helper: Stream from Azure OpenAI #
#######################################
def stream_azure_chat_completion(endpoint, headers, payload, print_stream=False):
    """
    Same streaming approach used in ask_func.py
    """
    with requests.post(endpoint, headers=headers, json=payload, stream=True) as response:
        response.raise_for_status()
        final_text = ""
        for line in response.iter_lines():
            if line:
                line_str = line.decode("utf-8", errors="ignore").strip()
                if line_str.startswith("data: "):
                    data_str = line_str[len("data: "):]
                    if data_str == "[DONE]":
                        break
                    try:
                        data_json = json.loads(data_str)
                        if (
                            "choices" in data_json
                            and data_json["choices"]
                            and "delta" in data_json["choices"][0]
                        ):
                            content_piece = data_json["choices"][0]["delta"].get("content", "")
                            if print_stream:
                                print(content_piece, end="", flush=True)
                            final_text += content_piece
                    except json.JSONDecodeError:
                        pass
        if print_stream:
            print()
    return final_text

########################################################
# 2) Call LLM to get PPT “instructions” from the AI    #
########################################################
def get_ppt_structure_from_ai(chat_history_str: str, instructions: str) -> str:
    """
    Calls Azure OpenAI with a specialized system prompt to produce
    textual instructions or structure for the slides. 
    We keep it simple: the LLM returns a textual description 
    of how the slides should look.
    """
    system_prompt = f"""
You are a PowerPoint presentation expert. Use the following information to make the slides.
Rules:
- Only use the following information to create the slides.
- Dont come up with anything outside of your scope in your slides.
- Your output will be utilized by the "python-pptx" to create the slides.


(The Information)
Conversation: 
{chat_history_str}

User_Instructions:
{instructions}
    """.strip()

    payload = {
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": "Please generate the PowerPoint structure now."}
        ],
        "max_tokens": 1200,
        "temperature": 0.7,
        "stream": True
    }

    headers = {
        "Content-Type": "application/json",
        "api-key": FAKE_LLM_API_KEY
    }

    # Stream from Azure:
    ppt_text = stream_azure_chat_completion(FAKE_LLM_ENDPOINT, headers, payload, print_stream=False)
    return ppt_text

######################################################################
# 3) Create PPT using python-pptx from the textual structure returned #
######################################################################
def create_ppt_from_structure(structure_text: str) -> BytesIO:
    """
    Given the AI output (structure_text) describing slides,
    parse it in a simple way, then build a .pptx using python-pptx.
    Returns a BytesIO of the generated PPT file.
    
    You can expand the parsing logic as needed.
    """
    # Create a blank Presentation
    prs = Presentation()
    
    # Simple parse: we look for lines that might indicate slides and content
    # For example, lines starting with "Slide " or "Title: ", "Content: "
    # This is a naive approach; you can refine as desired.
    
    lines = structure_text.splitlines()
    
    current_title = "Slide"
    current_body = []
    
    def add_slide(title, body_text):
        slide_layout = prs.slide_layouts[1]  # title+content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        slide.shapes.title.text = title
        
        # Set body
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = body_text
    
    for line in lines:
        # For illustration, let's assume the LLM might output something like:
        # Slide 1: Overview
        # Bullet: First point
        # Bullet: Second point
        # Slide 2: Another Title
        # ...
        # We'll do a basic parse:
        line_stripped = line.strip()
        
        # If we see "Slide" - that indicates start of a new slide
        slide_match = re.match(r"Slide\s*\d*:\s*(.+)", line_stripped, re.IGNORECASE)
        if slide_match:
            # If there's a previous slide's content, add it
            if current_body:
                add_slide(current_title, "\n".join(current_body))
            # Reset for new slide
            current_title = slide_match.group(1).strip()
            current_body = []
            continue
        
        # If line starts with "Bullet:" or similar, treat it as bullet point
        bullet_match = re.match(r"(Bullet|\-|\*):\s*(.+)", line_stripped, re.IGNORECASE)
        if bullet_match:
            current_body.append(f"• {bullet_match.group(2)}")
            continue
        
        # Otherwise, treat it as plain text
        if line_stripped:
            current_body.append(line_stripped)
    
    # Add the final slide if leftover content remains
    if current_body:
        add_slide(current_title, "\n".join(current_body))
    
    # Save to in-memory buffer
    ppt_buffer = BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    return ppt_buffer

##########################################################
# 4) Upload the .pptx to the same Azure storage + return #
##########################################################
def upload_ppt_to_azure(file_stream: BytesIO) -> str:
    """
    Uploads the PPTX file (file_stream) to the same Azure storage container
    used in your code. Returns a direct link (SAS URL) for download.
    """
    blob_service_client = BlobServiceClient(account_url=BLOB_ACCOUNT_URL, credential=BLOB_SAS_TOKEN)
    container_client = blob_service_client.get_container_client(BLOB_CONTAINER_NAME)
    
    timestamp_str = datetime.now().strftime("%Y%m%d_%H%M%S")
    ppt_filename = f"exported_ppt_{timestamp_str}.pptx"
    blob_name = BLOB_FOLDER_PATH + ppt_filename
    
    blob_client = container_client.get_blob_client(blob_name)
    
    file_stream.seek(0)
    blob_client.upload_blob(file_stream, overwrite=True)
    
    download_link = f"{BLOB_ACCOUNT_URL}/{BLOB_CONTAINER_NAME}/{blob_name}?{BLOB_SAS_TOKEN}"
    return download_link

##############################################
# 5) Main entry point to be called externally #
##############################################
def export_ppt(chat_history, ppt_instructions: str) -> str:
    """
    Main function to be called from elsewhere in your app:
    1) Collate the entire chat_history into one string
    2) Request the AI to produce slide structure
    3) Create a PPT using python-pptx
    4) Upload the PPT to Azure
    5) Return the direct download link
    """
    # Combine chat_history messages into a single string
    chat_history_str = "\n".join(chat_history)
    
    # 1. Get textual structure from AI
    structure_text = get_ppt_structure_from_ai(chat_history_str, ppt_instructions)
    
    # 2. Create the PPT using python-pptx
    ppt_buffer = create_ppt_from_structure(structure_text)
    
    # 3. Upload the PPT to Azure
    download_link = upload_ppt_to_azure(ppt_buffer)
    
    return download_link
