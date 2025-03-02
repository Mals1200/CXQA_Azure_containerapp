import requests
import json
import io
import threading
from datetime import datetime
from pptx import Presentation
from pptx.util import Pt
from azure.storage.blob import BlobServiceClient

def Call_PPT(latest_question, latest_answer, chat_history, instructions="No instructions provided"):
    """
    This function creates a PowerPoint based on the instructions, question/answer,
    and chat history. 
    If instructions == 'cancel', it cancels immediately. Otherwise, it calls Azure
    OpenAI to produce slide text, then builds the PPT and returns a download link.
    """

    # Exit the function if instructions are canceled:
    if instructions.lower() == "cancel":
        return "Operation canceled by the user."

    ##################################################
    # (A) CALL AZURE OPENAI TO GET SLIDE TEXT
    ##################################################
    chat_history_str = str(chat_history)

    ppt_prompt = f"""
You are a PowerPoint presentation expert. Use the following information to make the slides.
Rules:
- Only use the following information to create the slides.
- Don't come up with anything outside of your scope in your slides.
- Your output will be utilized by the "python-pptx" to create the slides (Keep in mind formatting).
- Make the output complete and ready for a presentation. **Only give the text for the slides**.
- Do not add any instructions or slide_numbers with the slides text.

(The Information)
- User_Instructions:
{instructions}

- Latest_Question:
{latest_question}

- Latest_Answer:
{latest_answer}

- Full Conversation:
{chat_history_str}
"""

    # FAKE Azure OpenAI values (example placeholders)
    LLM_ENDPOINT = (
        "https://cxqaazureaihub2358016269.openai.azure.com/"
        "openai/deployments/gpt-4o-3/chat/completions?api-version=2024-08-01-preview"
    )
    LLM_API_KEY = "xxxx"  # Your key

    system_message = "You are a helpful assistant that formats PowerPoint slides from user input."
    user_message = ppt_prompt

    payload = {
        "messages": [
            {"role": "system", "content": system_message},
            {"role": "user", "content": user_message}
        ],
        "max_tokens": 1000,
        "temperature": 0.7,
        "stream": True
    }
    headers = {
        "Content-Type": "application/json",
        "api-key": LLM_API_KEY
    }

    ppt_response = ""
    try:
        with requests.post(LLM_ENDPOINT, headers=headers, json=payload, stream=True) as response:
            response.raise_for_status()
            for line in response.iter_lines():
                if line:
                    line_str = line.decode("utf-8", errors="ignore").strip()
                    if line_str.startswith("data: "):
                        data_str = line_str[len("data: "):]
                        if data_str == "[DONE]":
                            break
                        try:
                            data_json = json.loads(data_str)
                            if (
                                "choices" in data_json
                                and data_json["choices"]
                                and "delta" in data_json["choices"][0]
                            ):
                                content_piece = data_json["choices"][0]["delta"].get("content", "")
                                ppt_response += content_piece
                        except json.JSONDecodeError:
                            pass
    except Exception as e:
        ppt_response = f"An error occurred while creating the PowerPoint slides: {e}"

    slides_text = ppt_response.strip()
    if not slides_text or "An error occurred" in slides_text:
        return f"No valid slide text returned:\n{slides_text}"

    ##################################################
    # (B) PARSE TEXT INTO SLIDES & CREATE A .PPTX
    ##################################################
    raw_slides = [s.strip() for s in slides_text.split("\n\n") if s.strip()]

    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title & Content

    for raw_slide in raw_slides:
        lines = raw_slide.split("\n")
        title_text = lines[0]
        bullet_lines = lines[1:]

        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title_text

        if bullet_lines:
            body_tf = slide.placeholders[1].text_frame
            body_tf.clear()
            for bullet_item in bullet_lines:
                p = body_tf.add_paragraph()
                p.text = bullet_item
                p.font.size = Pt(18)  # Example bullet font size

    ppt_buffer = io.BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)

    ##################################################
    # (C) UPLOAD TO AZURE BLOB STORAGE
    ##################################################
    account_url = "https://cxqaazureaihub8779474245.blob.core.windows.net"
    sas_token = (
        "sv=2022-11-02&ss=bfqt&srt=sco&sp=rwdlacupiytfx&"
        "se=2030-11-21T02:02:26Z&st=2024-11-20T18:02:26Z&"
        "spr=https&sig=xxxx"
    )
    container_name = "5d74a98c-1fc6-4567-8545-2632b489bd0b-azureml-blobstore"

    blob_service_client = BlobServiceClient(account_url=account_url, credential=sas_token)
    container_client = blob_service_client.get_container_client(container_name)

    ppt_filename = f"slides_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    blob_client = container_client.get_blob_client(ppt_filename)

    blob_client.upload_blob(ppt_buffer, overwrite=True)

    download_link = f"{account_url}/{container_name}/{ppt_filename}?{sas_token}"

    ##################################################
    # (D) SCHEDULE AUTO-DELETE AFTER 5 MINUTES
    ##################################################
    def delete_blob_after_5():
        try:
            blob_client.delete_blob()
        except Exception:
            pass

    timer = threading.Timer(300, delete_blob_after_5)  # 300s = 5 minutes
    timer.start()

    # (E) Return the link
    return download_link
