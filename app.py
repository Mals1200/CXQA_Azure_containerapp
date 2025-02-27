import os
import asyncio
from flask import Flask, request, jsonify, Response
import requests

# Bot Framework imports
from botbuilder.core import BotFrameworkAdapter, BotFrameworkAdapterSettings, TurnContext
from botbuilder.schema import Activity

# Import your existing logic from ask_func.py
from ask_func import Ask_Question

# Additional libraries for PPT
from pptx import Presentation
import uuid
from azure.storage.blob import BlobServiceClient

app = Flask(__name__)

# Read Bot credentials from environment variables
MICROSOFT_APP_ID = os.environ.get("MICROSOFT_APP_ID", "")
MICROSOFT_APP_PASSWORD = os.environ.get("MICROSOFT_APP_PASSWORD", "")

# Create settings & adapter for the Bot
adapter_settings = BotFrameworkAdapterSettings(MICROSOFT_APP_ID, MICROSOFT_APP_PASSWORD)
adapter = BotFrameworkAdapter(adapter_settings)

###########################################################
# Conversation State
# - We now store multiple Q&A answers in a dictionary
#   { answer_id: <answer text> }
#
# conversation_histories[conversation_id] = [
#    ... chat messages ...
# ]
# conversation_histories[conversation_id + "_answers"] = {
#    <uuid>: <answer_text>,
#    ...
# }
###########################################################
conversation_histories = {}

@app.route('/', methods=['GET'])
def home():
    return jsonify({'message': 'API is running!'}), 200

@app.route('/ask', methods=['POST'])
def ask():
    data = request.get_json()
    if not data or 'question' not in data:
        return jsonify({'error': 'Invalid request, "question" field is required.'}), 400
    question = data['question']
    # For direct usage, we still rely on the global chat_history in ask_func.
    answer = Ask_Question(question)
    return jsonify({'answer': answer})

@app.route("/api/messages", methods=["POST"])
def messages():
    if "application/json" not in request.headers.get("Content-Type", ""):
        return Response(status=415)
    # Deserialize incoming Activity
    body = request.json
    activity = Activity().deserialize(body)
    # Get the Authorization header (for Bot Framework auth)
    auth_header = request.headers.get("Authorization", "")
    loop = asyncio.new_event_loop()
    try:
        loop.run_until_complete(adapter.process_activity(activity, auth_header, _bot_logic))
    finally:
        loop.close()
    return Response(status=200)

async def _bot_logic(turn_context: TurnContext):
    # Unique conversation ID
    conversation_id = turn_context.activity.conversation.id

    # 1) Ensure we have a chat history list for this conversation
    if conversation_id not in conversation_histories:
        conversation_histories[conversation_id] = []

    # 2) Also ensure we have an "answers" dict to store multiple answers
    answers_key = conversation_id + "_answers"
    if answers_key not in conversation_histories:
        conversation_histories[answers_key] = {}

    # Hook up ask_func's chat_history to this conversation's messages
    import ask_func
    ask_func.chat_history = conversation_histories[conversation_id]

    # Check if the user clicked an Action.Submit button
    if turn_context.activity.value and "action" in turn_context.activity.value:
        action = turn_context.activity.value["action"]

        if action == "export":
            # The user wants to export a PPT from a previous answer
            answer_id = turn_context.activity.value.get("answer_id", "")
            # Retrieve that specific answer text from the answers dict
            text_for_ppt = conversation_histories[answers_key].get(answer_id, "")

            if not text_for_ppt:
                # If there's no matching answer, we can't export
                await turn_context.send_activity("No previous answer found to export.")
                return

            # Otherwise, generate PPT
            ppt_filename = generate_and_upload_ppt(text_for_ppt)
            if ppt_filename:
                await turn_context.send_activity(
                    f"I've generated your PPT! You can download it here:\n{ppt_filename}"
                )
            else:
                await turn_context.send_activity("Sorry, I couldn't create the PPT file.")
        else:
            await turn_context.send_activity("Unknown action received.")
        return

    # Otherwise, it's a normal user message: we ask the question
    user_message = turn_context.activity.text or ""
    answer_text = Ask_Question(user_message)

    # Update the conversation's chat messages
    conversation_histories[conversation_id] = ask_func.chat_history

    # 3) Store the new answer in the answers dict using a unique ID
    new_answer_id = str(uuid.uuid4())
    conversation_histories[answers_key][new_answer_id] = answer_text

    # If the answer has "Source:" content, we show an Adaptive Card. 
    # Or, you can always show an Adaptive Card for every answer, if you want.
    if "\n\nSource:" in answer_text:
        # Split main answer vs. source
        parts = answer_text.split("\n\nSource:", 1)
        main_answer = parts[0].strip()
        source_details = "Source:" + parts[1].strip()

        # Build a card with a Show Source button & an Export PPT button
        adaptive_card = {
            "type": "AdaptiveCard",
            "body": [
                {"type": "TextBlock", "text": main_answer, "wrap": True},
                {
                    "type": "TextBlock",
                    "text": source_details,
                    "wrap": True,
                    "id": "sourceBlock",
                    "isVisible": False
                }
            ],
            "actions": [
                {
                    "type": "Action.ToggleVisibility",
                    "title": "Show Source",
                    "targetElements": ["sourceBlock"]
                },
                # Notice we pass "action": "export" AND the "answer_id" so we can find the right answer
                {
                    "type": "Action.Submit",
                    "title": "Export PPT",
                    "data": {
                        "action": "export",
                        "answer_id": new_answer_id  # NEW
                    }
                }
            ],
            "$schema": "http://adaptivecards.io/schemas/adaptive-card.json",
            "version": "1.2"
        }

        message = Activity(
            type="message",
            attachments=[{
                "contentType": "application/vnd.microsoft.card.adaptive",
                "content": adaptive_card
            }]
        )
        await turn_context.send_activity(message)
    else:
        # If there's no "Source:" line, we can still show a card with "Export PPT" if you prefer.
        # Or just plain text. For demonstration, let's show a simpler card with an export button:

        adaptive_card = {
            "type": "AdaptiveCard",
            "body": [
                {"type": "TextBlock", "text": answer_text, "wrap": True}
            ],
            "actions": [
                {
                    "type": "Action.Submit",
                    "title": "Export PPT",
                    "data": {
                        "action": "export",
                        "answer_id": new_answer_id  # NEW
                    }
                }
            ],
            "$schema": "http://adaptivecards.io/schemas/adaptive-card.json",
            "version": "1.2"
        }

        message = Activity(
            type="message",
            attachments=[{
                "contentType": "application/vnd.microsoft.card.adaptive",
                "content": adaptive_card
            }]
        )
        await turn_context.send_activity(message)

##############################################
# PPT Generation Helper
##############################################
def generate_and_upload_ppt(text_for_ppt):
    """
    1) Generate a local PPTX file from the answer text.
    2) Upload it to Azure Blob Storage (or any location).
    3) Return a direct link or None if error.
    """
    try:
        # 1) Generate PPT using python-pptx
        prs = Presentation()
        # Slide 1: Title & Subtitle
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Your Generated PPT"
        slide.placeholders[1].text = "Subtitle from GPT"

        # Slide 2: bullet points
        bullet_slide_layout = prs.slide_layouts[1]
        slide2 = prs.slides.add_slide(bullet_slide_layout)
        slide2.shapes.title.text = "Main Answer"
        body_shape = slide2.shapes.placeholders[1]
        tf = body_shape.text_frame

        for line in text_for_ppt.split("\n"):
            if line.strip():
                p = tf.add_paragraph()
                p.text = line.strip()

        # Save to a temporary file
        filename = f"ppt_{uuid.uuid4()}.pptx"
        prs.save(filename)

        # 2) Upload to Azure Blob
        #    (Fake placeholders—replace with real credentials)
        account_url = "https://cxqaazureaihub8779474245.blob.core.windows.net"
        sas_token = (
            "sv=2022-11-02&ss=bfqt&srt=sco&sp=rwdlacupiytfx&"
            "se=2030-11-21T02:02:26Z&st=2024-11-20T18:02:26Z&"
            "spr=https&sig=YfZEUMeqiuBiG7le2JfaaZf%2FW6t8ZW75yCsFM6nUmUw%3D"
        )
        container_name = "5d74a98c-1fc6-4567-8545-2632b489bd0b-azureml-blobstore"
        target_blob_name = f"UI/ppt_exports/{filename}"

        blob_service_client = BlobServiceClient(account_url=account_url, credential=sas_token)
        container_client = blob_service_client.get_container_client(container_name)

        with open(filename, "rb") as data:
            container_client.upload_blob(name=target_blob_name, data=data, overwrite=True)

        # Create a direct link to the PPT
        ppt_url = f"{account_url}/{container_name}/{target_blob_name}?{sas_token}"

        # Remove local file
        import os
        os.remove(filename)

        return ppt_url
    except Exception as e:
        print(f"Error generating/uploading PPT: {e}")
        return None

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=80)
